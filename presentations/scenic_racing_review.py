"""SD-17: 60-min presentation deck generator for the Scenic racing domain talk.

Audience: core Scenic developers (academic). Style: academic-conference
(white, Calibri). Three-section narrative: scenarios -> racing library ->
dSPACE integration, then constructive feedback to Scenic core.

Run:
    python presentations/scenic_racing_review.py

Output:
    presentations/scenic_racing_review.pptx

The .pptx imports cleanly into Google Slides via Drive upload.

Code-anchor lines hyperlink to the canonical branch on GitHub
(https://github.com/Lawrence1024/Scenic/tree/cosim) so the speaker
can click during the live talk to jump to the file:line.

Re-runs are idempotent (overwrites the .pptx).
"""

from __future__ import annotations

import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------------
# Code-anchor -> GitHub URL
# ---------------------------------------------------------------------------

GITHUB_BASE = "https://github.com/Lawrence1024/Scenic"
BRANCH = "cosim"

# Roots we recognize as "this is a repo path" (vs. arbitrary text).
_ROOT_RE = r"(?:src|examples|docs|tests|tools|assets|presentations)"


def _anchor_url(text: str) -> Optional[str]:
    """Build a GitHub URL for the FIRST repo path mentioned in `text`.

    Recognizes:
        src/.../file.ext        -> blob URL
        src/.../file.ext:42     -> blob URL with #L42 fragment
        src/.../file.ext:42-99  -> blob URL with #L42-L99 fragment
        src/.../somedir/        -> tree URL (directory)

    Returns None if no recognizable repo path is found.
    """
    m = re.search(
        rf"({_ROOT_RE}/[\w./-]+\.\w+)(?::(\d+)(?:-(\d+))?)?",
        text,
    )
    if m:
        path, l1, l2 = m.group(1), m.group(2), m.group(3)
        url = f"{GITHUB_BASE}/blob/{BRANCH}/{path}"
        if l1 and l2:
            url += f"#L{l1}-L{l2}"
        elif l1:
            url += f"#L{l1}"
        return url
    m = re.search(rf"({_ROOT_RE}/[\w./-]+?)/?(?:\s|$|\()", text + " ")
    if m:
        path = m.group(1).rstrip("/")
        return f"{GITHUB_BASE}/tree/{BRANCH}/{path}"
    return None


# ---------------------------------------------------------------------------
# Style constants -- academic-conference (light, sans-serif)
# ---------------------------------------------------------------------------

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"
CODE_FONT = "Cascadia Code"

TITLE_SIZE = Pt(32)
BODY_SIZE = Pt(18)
SUB_SIZE = Pt(16)
CODE_SIZE = Pt(14)
FOOTER_SIZE = Pt(11)
MISSING_SIZE = Pt(13)
NOTE_SIZE = Pt(12)

COLOR_TITLE = RGBColor(0x1F, 0x3A, 0x6E)
COLOR_BODY = RGBColor(0x20, 0x20, 0x20)
COLOR_CODE = RGBColor(0x50, 0x50, 0x50)
COLOR_LINK = RGBColor(0x05, 0x63, 0xC1)        # standard hyperlink blue
COLOR_FOOTER = RGBColor(0x80, 0x80, 0x80)
COLOR_MISSING = RGBColor(0xC4, 0x55, 0x00)
COLOR_VIDEO_BORDER = RGBColor(0xB0, 0xB0, 0xB0)
COLOR_VIDEO_FILL = RGBColor(0xF6, 0xF6, 0xF6)

TITLE_BOX = (Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.85))
BODY_BOX = (Inches(0.7), Inches(1.3), Inches(8.2), Inches(4.6))
BODY_BOX_FULL = (Inches(0.7), Inches(1.3), Inches(12.0), Inches(4.6))
VIDEO_BOX = (Inches(9.2), Inches(1.5), Inches(3.6), Inches(2.4))
CODE_BOX = (Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.55))
MISSING_BOX = (Inches(0.7), Inches(6.65), Inches(12.0), Inches(0.45))
FOOTER_BOX = (Inches(0.5), Inches(7.10), Inches(12.3), Inches(0.35))


# ---------------------------------------------------------------------------
# Slide data model
# ---------------------------------------------------------------------------

@dataclass
class Slide:
    n: int
    section: str
    title: str
    bullets: List[str] = field(default_factory=list)
    code_anchor: str = ""
    video_spec: Optional[str] = None
    missing_marker: Optional[str] = None
    speaker_note: str = ""
    # Side-by-side code blocks. Each is (label, code_text). Set both for
    # a comparison layout (e.g. "today's code" vs "proposed Scenic syntax"
    # on the Ask-2 code-comparison slide). When set, bullets are skipped
    # for that slide -- the body box is replaced by two code panels.
    code_left: Optional[tuple] = None
    code_right: Optional[tuple] = None
    # Set True to render the planner FSM diagram in the right panel
    # of the slide (where video placeholders normally go). Used by
    # the FSM-detail slide to show the 5-state machine visually
    # alongside bullets.
    fsm_diagram: bool = False
    # Set True for topic-divider slides: title rendered large and
    # centered, optional one-line subtitle from bullets[0]. No code
    # anchor, no missing-marker. Used as section breaks so the
    # audience sees a clean transition between major topics.
    is_divider: bool = False


# ---------------------------------------------------------------------------
# Section names ("Topic" rather than "Act" per audience-friendliness)
# ---------------------------------------------------------------------------

SECTION_OPEN = "Open"
SECTION_T1 = "Topic 1: Scenarios"
SECTION_T2A = "Topic 2a: Geometry"
SECTION_T2B = "Topic 2b: Behaviors + MPC"
SECTION_T2C = "Topic 2c: Falsification"
SECTION_T3 = "Topic 3: dSPACE"
SECTION_FB = "Feedback"
SECTION_CLOSE = "Close"


# ---------------------------------------------------------------------------
# All 39 slides
# ---------------------------------------------------------------------------

SLIDES: List[Slide] = [
    # ---- Cover (divider) ----
    Slide(
        n=1, section=SECTION_OPEN, is_divider=True,
        title="Racing on Scenic",
        bullets=[
            "A Domain Built On Top of Driving",
            "By: Lawrence Shieh",
        ],
        speaker_note="",
    ),

    # ---- Open (intro slides) ----
    Slide(
        n=2, section=SECTION_OPEN,
        title="Racing on Scenic: A Domain Built On Top of Driving",
        bullets=[
            "Lawrence Shieh (5th Year MS Student)",
            "Presentation Focus: Scenic's class / region / behavior model carried us through ~67 mph racing -- with two specific gaps we'll itemize at the end",
            "3-topic roadmap: scenarios -> racing library -> dSPACE bridge -> feedback",
            "Code-tour interleaved -- code anchors are clickable; open the repo now",
        ],
        code_anchor="src/scenic/domains/racing/  (root)",
        video_spec="8s loop of IAC car on F-bank as visual hook",
        speaker_note="",
    ),
    Slide(
        n=3, section=SECTION_OPEN,
        title="Why Racing Stresses Scenic Differently",
        bullets=[
            "Closing speeds ~67 mph (30 m/s) vs. urban driving's ~22 mph (10 m/s)",
            "Oval / road-course geometry, not grid lanes",
            "Competitive intent (overtake, block) needs strategy, not just maneuvers",
            "Falsification has hard real-time cost (dSPACE wall-clock)",
        ],
        code_anchor="src/scenic/domains/racing/model.scenic:125-196  (RacingCar delta)",
        speaker_note="",
    ),
    Slide(
        n=4, section=SECTION_OPEN,
        title="Roadmap & Today's Deliverable",
        bullets=[
            "Topic 1: scenario banks -- F-bank (regression) + falsifiable (exploration)",
            "Topic 2: racing library internals -- geometry, behaviors, MPC, planner, falsifier",
            "Topic 3: dSPACE bridge -- light architectural sketch",
            "Closing: two concrete asks with API sketches and implementation paths",
            "Time: ~30 min slides + ~25 min live code + ~5 min Q&A",
        ],
        speaker_note="",
    ),

    # ---- Topic 1: Scenarios (divider + 4 slides) ----
    Slide(
        n=5, section=SECTION_T1, is_divider=True,
        title="Topic 1",
        bullets=["Scenarios -- two banks, one shared library"],
        speaker_note="",
    ),
    Slide(
        n=6, section=SECTION_T1,
        title="Two Scenario Banks, One Shared Library",
        bullets=[
            "F-bank -- regression suite, 13 scenarios (F0-F12), hardcoded positions",
            "Falsifiable -- one .scenic file (S1_falsify), distribution-driven",
            "Same .scenic file shape across both; only the parameter type differs",
            "F-bank locks specific failure modes; falsifiable explores parameter space",
        ],
        code_anchor="examples/racing/  (browse f_shared/ and falsifiable/)",
        speaker_note="",
    ),
    Slide(
        n=7, section=SECTION_T1,
        title="F-Bank Walk: F0-F5 (Solo and Reactive Fellows)",
        bullets=[
            "F0  ego-only baseline (no fellow)",
            "F1  fellow behind, no interaction",
            "F2  ahead, slower -- the overtake setup",
            "F3L / F3R  side-blocked (left or right TTL)",
            "F4  sudden stop -- SuddenStopInterval behavior",
            "F5  swerve -- SwerveOutOfControl behavior",
        ],
        code_anchor="examples/racing/f_shared/F4_fellow_ahead_sudden_stop.scenic",
        video_spec="dSPACE viewer F4 emergency-brake event, ~20 s",
        speaker_note="",
    ),
    Slide(
        n=8, section=SECTION_T1,
        title="F-Bank Walk: F6-F12 (Geometric and Strategic)",
        bullets=[
            "F6 / F7  strict left / right occupied (passing-side forced)",
            "F8  corner + ahead -- by-design conservatism",
            "F9  stationary roadside obstacle (speed=0 given to fellow behavior)",
            "F10 / F11 / F12  corner-entry variants (left/right/stop)",
            "Each scenario regresses against a specific failure mode",
        ],
        code_anchor="examples/racing/f_shared/F9_fellow_stationary_roadside_obstacle.scenic",
        video_spec="F9 stationary obstacle clip showing ego cruising past with speed=0 fellow, ~20 s",
        speaker_note="",
    ),
    Slide(
        n=9, section=SECTION_T1,
        title="Falsifiable Bank: S1_falsify.scenic",
        bullets=[
            "One file: examples/racing/falsifiable/S1_falsify.scenic (~50 LOC)",
            "Key line: gap_m = VerifaiRange(20, 60)  -- the only special syntax",
            "Everything else is plain Scenic (model, ego, opponent, behaviors)",
            "Sampler chosen at runner-time: --sampler {halton, ce, bo, random}",
            "Halton -> uniform coverage; CE -> active falsification (we'll see CE results in Topic 2c)",
        ],
        code_anchor="examples/racing/falsifiable/S1_falsify.scenic:41",
        speaker_note="",
    ),

    # ---- Topic 2a: Geometry (divider + 6 slides) ----
    Slide(
        n=10, section=SECTION_T2A, is_divider=True,
        title="Topic 2a",
        bullets=["Geometry -- RacingCar, RacingTrack, regions, TTL"],
        speaker_note="",
    ),
    Slide(
        n=11, section=SECTION_T2A,
        title="Inheritance Chain: Vehicle -> Car -> RacingCar",
        bullets=[
            "Vehicle  driving/model.scenic:268",
            "Car      driving/model.scenic:298",
            "RacingCar  racing/model.scenic:125",
            "RacingCar adds: raceNumber, IAC dims (1.93 x 4.88 m / 76 x 192 in), maxSpeed=~67 mph, ttlFolder, ttlFileName",
            "Total delta: ~70 lines",
        ],
        code_anchor="src/scenic/domains/racing/model.scenic:125-196",
        speaker_note="",
    ),
    Slide(
        n=12, section=SECTION_T2A,
        title="What RacingCar Adds: Field-by-Field Explanation",
        bullets=[
            "Identity: raceNumber, team, carType  -- needed by dSPACE traffic + eval logs",
            "Geometry: width=1.93 m, length=4.88 m  -- IAC Dallara AV-24; drives the OBB collision check",
            "Performance envelope: maxSpeed=~67 mph, accel=8 m/s^2, braking=-12 m/s^2  -- feeds planner + MPC speed cap",
            "State (aspirational): fuelLevel, tireWear  -- extension points for endurance scenarios; not used today",
            "Controller knobs: waypointTolerance=2.5 m, controllerAggressiveness=0.5  -- MPC tuning",
            "TTL routing: ttlFolder, ttlFileName  -- per-vehicle racing-line selection",
        ],
        code_anchor="src/scenic/domains/racing/model.scenic:125-196",
        speaker_note="",
    ),
    Slide(
        n=13, section=SECTION_T2A,
        title="The Default Position Trick",
        bullets=[
            "position: new Point on ttlRegion(self.ttlFileName)",
            "Self-referential default: depends on this car's ttlFileName property",
            "Composes with inherited parentOrientation = (roadDirection at position) + roadDeviation",
            "Result: `new RacingCar` with no specifiers = track-aligned car at a valid TTL position",
            "Explicit specifiers (`new RacingCar on mainTrack`) still override",
        ],
        code_anchor="src/scenic/domains/racing/model.scenic:164",
        speaker_note="",
    ),
    Slide(
        n=14, section=SECTION_T2A,
        title="RacingTrack: Wrapping Scenic's Network",
        bullets=[
            "Scenic's driving domain provides Network -- parsed from XODR (OpenDRIVE)",
            "Network gives us roads, lanes, intersections, centerlines, lane edges",
            "RacingTrack (segments/tracks.py) is a Python class that wraps that Network",
            "It identifies main-loop roads vs. pit-lane roads via topology + name patterns",
            "Exposes: track.network, track._mainRacingRoads, track._pitRoads",
        ],
        code_anchor="src/scenic/domains/racing/segments/tracks.py:807-852",
        speaker_note="",
    ),
    Slide(
        n=15, section=SECTION_T2A,
        title="Region Composition: mainTrack, pitTrack, ttlRegion",
        bullets=[
            "mainTrack, pitTrack -- buffered polygons around road centerlines (XODR-native)",
            "ttlRegion(file) -- buffered polygon around a TTL CSV centerline (per-line)",
            "Ego placement: explicit `on mainTrack` OR per-vehicle default `on ttlRegion(self.ttlFileName)`",
            "TTL CSV format: x, y, z columns (3D centerline waypoints)",
            "Multiple TTLs per track: optimal, left, right, pit",
        ],
        code_anchor="src/scenic/domains/racing/segments/track_regions.py:275-352",
        missing_marker="Missing from Scenic: no first-class Frenet/RD region.",
        speaker_note="",
    ),
    Slide(
        n=16, section=SECTION_T2A,
        title="The Frenet Gap (Foreshadow)",
        bullets=[
            "Every racing scenario reasons in (s, t) coordinates",
            "Scenic gives us (x, y) regions; conversion is per-simulator",
            "We built a parallel RD frame in Python (~30 LOC, dSPACE-only)",
            "Surface today: dSPACE specifier `_racing_st_offset` (Topic 3)",
            "Proposal at close: Scenic core could absorb this",
        ],
        code_anchor="src/scenic/simulators/dspace/modeldesk/placement.py:26-58",
        speaker_note="",
    ),

    # ---- Topic 2b: Behaviors + MPC (divider + 8 slides) ----
    Slide(
        n=17, section=SECTION_T2B, is_divider=True,
        title="Topic 2b",
        bullets=["Behaviors + MPC -- the star behavior, planner, FSM, safety"],
        speaker_note="",
    ),
    Slide(
        n=18, section=SECTION_T2B,
        title="The Star Behavior: FollowRacingLineMPCBehavior",
        bullets=[
            "Defined at behaviors.scenic:310 (signature with ~13 kwargs)",
            "Main loop spans lines 554-2840 (yes, ~2300 lines)",
            "Called from every ego scenario in F-bank and falsifiable",
            "Composes with `try-interrupt` for events (collisions, mode switches)",
        ],
        code_anchor="src/scenic/domains/racing/behaviors.scenic:310",
        speaker_note="",
    ),
    Slide(
        n=19, section=SECTION_T2B,
        title="The 8-Phase Tick (Survey)",
        bullets=[
            "1. State read (sim/dSPACE readback)",
            "2. Waypoint progress (advance index on current TTL -> gives ego arc-length `s`)",
            "3. Tactical planner (uses `s`; picks TTL + speed cap)",
            "4. MPC reference build (sample chosen-TTL waypoints over horizon)",
            "5. Lateral MPC -> steering",
            "6. Longitudinal MPC -> throttle / brake",
            "7. Gear logic (rule-based, post-MPC)",
            "8. Safety gates -> action emit",
        ],
        code_anchor="src/scenic/domains/racing/behaviors.scenic:554-2840",
        speaker_note="",
    ),
    Slide(
        n=20, section=SECTION_T2B,
        title="Lateral MPC (Deep Into Params)",
        bullets=[
            "State: [e_y, e_psi, delta, s]  (cross-track err, heading err, steer, progress)",
            "Control: steering rate (rad/s)  -- integrated to delta",
            "Solver: OSQP (sparse QP); ~5-10 ms per tick",
            "Horizon: 35 steps x 0.05 s = 1.75 s  (the LOCAL tracking horizon)",
            "Reference: TTL spline -> (x_ref, y_ref, psi_ref, kappa_ref) per horizon step",
            "Cost: contouring (e_y, e_psi) + smoothness (delta_rate^2) + progress reward",
        ],
        code_anchor="src/scenic/domains/racing/mpc/reference_builder.py:14-200",
        speaker_note="",
    ),
    Slide(
        n=21, section=SECTION_T2B,
        title="Longitudinal MPC + Gear Logic (Deep Into Params)",
        bullets=[
            "State: [v, a]  -- velocity, acceleration",
            "Control: [throttle, brake]  in [0, 1]^2",
            "Same 1.75 s horizon as lateral",
            "v_target from TTL line (curvature-aware speed cap: v_max = sqrt(a_y_max / kappa))",
            "Coupling: longitudinal v feeds lateral kappa lookup",
            "Gear logic: rule-based AFTER MPC (proactive downshift before high-curvature)",
        ],
        code_anchor="src/scenic/domains/racing/behaviors.scenic:554-2840",
        speaker_note="",
    ),
    Slide(
        n=22, section=SECTION_T2B,
        title="Tactical Planner: The Strategic Layer",
        bullets=[
            "Picks one of 4 strategies: stay_optimal, follow_fellow, pass_left, pass_right",
            "Trajectory prediction: fellow over 10 s horizon (constant-velocity model)",
            "Strategy simulator: rolls out each candidate; ranks by progress + min-clearance",
            "Output: planner mode + chosen TTL + speed cap",
            "Then hands the chosen mode to the executor state machine (next slide)",
        ],
        code_anchor="src/scenic/domains/racing/tactical_planner.py:350-430",
        speaker_note="",
    ),
    Slide(
        n=23, section=SECTION_T2B,
        title="The Executor State Machine -- Hand-Rolled FSM",
        bullets=[
            "5 logical states (COMMIT/HOLD have L/R variants -> 7 constants)",
            "Initial: FREE_RUN; no terminal state",
            "Storage: `TacticalPlannerState.mode: str = FREE_RUN`",
            "Transitions: chained `if state.mode == X` branches",
            "Entry actions: inlined Python beside each `state.mode = X`",
            "Output: (mode, ttl, speed_cap, decision_reason) per tick",
        ],
        code_anchor="src/scenic/domains/racing/tactical_planner.py:330-396",
        missing_marker="Missing from Scenic: no native state-machine construct for behaviors.",
        fsm_diagram=True,
        speaker_note="",
    ),
    Slide(
        n=24, section=SECTION_T2B,
        title="Safety Gate: Predicted-Collision Override",
        bullets=[
            "path_collision_predicted(ego_path, opp_trajectory, horizon_s=1.5)",
            "Sample dt 0.1 s -> 15 samples per check",
            "Min clearance threshold: 0.5 m (hard overlap)",
            "Debounced (2 consecutive breaches) to handle numerical noise",
            "Independent of strategy layer -- can override even mid-COMMIT",
        ],
        code_anchor="src/scenic/domains/racing/assessment/pass_geometry.py:194-240",
        speaker_note="",
    ),
    Slide(
        n=25, section=SECTION_T2B,
        title="Fellow Behaviors (Brief)",
        bullets=[
            "FellowFollowTTLGeometricBehavior  -- thin Scenic wrapper at behaviors.scenic:183-207",
            "Real logic: compute_follow_ttl_geometric_plant_command  in fellow/commands.py:688-700",
            "Per tick: constant v_kmh + delta(s) feedforward from TTL waypoint table; no PID, no MPC",
            "FellowSuddenStopIntervalBehavior  -- periodic stops (F4); FellowSwerveOutOfControl  (F5)",
            "Adversarial in EFFECT (causes trouble for ego) but not LEARNING -- scripted disturbances",
        ],
        code_anchor="src/scenic/domains/racing/fellow/commands.py:688-700",
        speaker_note="",
    ),

    # ---- Topic 2c: Falsification (divider + 6 slides) ----
    Slide(
        n=26, section=SECTION_T2C, is_divider=True,
        title="Topic 2c",
        bullets=["Falsification -- the runner, monitors, and the 50-sample result"],
        speaker_note="",
    ),
    Slide(
        n=27, section=SECTION_T2C,
        title="The Falsifier Loop -- One Runner, Pluggable Sampler",
        bullets=[
            "verifai_runner.py -- single in-process runner; handles every sampling mode",
            "Compiles scenario once; cosim bridge stays warm across ALL samples",
            "~30 s VEOS init paid once per campaign, regardless of sample count",
            "Sampler swap via --sampler {halton, ce, bo, random}; same .scenic, same logs",
            "Halton = quasi-random coverage smoke; CE = active falsification",
        ],
        code_anchor="src/scenic/domains/racing/benchmarks/verifai_runner.py",
        speaker_note="",
    ),
    Slide(
        n=28, section=SECTION_T2C,
        title="Inside the Loop -- main() Walkthrough",
        bullets=[
            "Setup (once per campaign):",
            "  scenarioFromFile(file, params={'verifaiSamplerType': args.sampler})  -- compile, inject sampler",
            "  simulator = scenario.getSimulator()  -- spawns cosim bridge on first .simulate()",
            "Per sample (count iterations):",
            "  scene = scenario.generate(feedback=last_rho)  -- VerifAI's CE picks next gap",
            "  simulator.simulate(scene)  -- stdout teed to logs/sample_NNN.log",
            "  metrics = parse_sample(log_path)  -- regex-extracts SampleMetrics from STDOUT",
            "  rho = monitors.RESOLVE[args.monitor](metrics)  ;  feedback = rho",
            "Caveat: monitors read PARSED STDOUT, not in-memory state -- a logging change would silently break this",
        ],
        code_anchor="src/scenic/domains/racing/benchmarks/verifai_runner.py:352",
        speaker_note="",
    ),
    Slide(
        n=29, section=SECTION_T2C,
        title="Monitors -- Robustness Predicates Over a Run",
        bullets=[
            "collision  -- bbox_gap_m_min: minimum bounding-box gap (meters); 0 = touch, <0 = overlap",
            "track  -- track_clearance_m: signed distance to nearest track edge (meters)",
            "off-track / brake / overtake-success  -- semi-continuous count signals",
            "safety = min(collision, track)  -- two-spec composite (the one we used in the campaign)",
            "min  -- four-way composite over all safety + planner signals",
            "all  -- multi-objective tuple for `mab` (multi-armed bandit) samplers",
        ],
        code_anchor="src/scenic/domains/racing/benchmarks/monitors.py",
        speaker_note="",
    ),
    Slide(
        n=30, section=SECTION_T2C,
        title="The CE Run That Broke the Planner",
        bullets=[
            "50-sample cross-entropy run (--monitor safety, --time 3000)",
            "29 / 50 collisions  (58%)",
            "Collisions span gap in [21.7, 52.3] m -- broad failure, NOT a narrow cluster",
            "6 successful overtakes  (all via commit_pass_right)",
        ],
        code_anchor="src/scenic/domains/racing/benchmarks/results/verifai_20260428_052048/",
        video_spec="30 s replay of a representative collision (e.g. sample 11 gap=22.0m, or sample 17 gap=30.6m for a long-commit case)",
        speaker_note="",
    ),
    Slide(
        n=31, section=SECTION_T2C,
        title="What's Failing: Hypothesis from 50 Samples",
        bullets=[
            "Striking asymmetry: 28 / 29 collisions involve LEFT-pass attempts only (1 mixed, 0 right-only)",
            "Pattern in left-pass collisions: 50-500 ticks of commit_pass_left_hold + 60-120 aborts -> contact",
            "Meanwhile: all 6 successful overtakes happened on COMMIT_PASS_RIGHT",
            "Open: TTL geometry on the left line? MPC lateral limits? Strategy preferring left over right?",
            "Off-track instrumentation invalid this run: ~10 m frame residual + 'd_in=0 = OUT' semantics",
        ],
        code_anchor="src/scenic/domains/racing/tactical_planner.py:350-430",
        speaker_note="",
    ),

    Slide(
        n=32, section=SECTION_T2C,
        title="The Pipeline Is the Deliverable",
        bullets=[
            "58% collision rate is a Smart-Ego result, NOT a pipeline-quality result",
            "Falsifier is ego-agnostic: ego is just a Scenic behavior assigned in the .scenic file",
            "Swap `FollowRacingLineMPCBehavior` for the ART stack -> falsify ART under the same scenarios",
            "Same .scenic, same monitors, same VerifAI sampler -- only `ego.behavior` changes",
            "The 28/29 left-only finding is ABOUT Smart Ego; another stack will surface its own asymmetries",
        ],
        code_anchor="examples/racing/falsifiable/S1_falsify.scenic:55",
        speaker_note="",
    ),

    # ---- Topic 3: dSPACE (divider + 5 slides) ----
    Slide(
        n=33, section=SECTION_T3, is_divider=True,
        title="Topic 3",
        bullets=["dSPACE -- VEOS, the cosim bridge, MAPort, _racing_st_offset"],
        speaker_note="",
    ),
    Slide(
        n=34, section=SECTION_T3,
        title="Why dSPACE? VEOS in One Slide",
        bullets=[
            "VEOS = dSPACE's HIL-compatible offline simulator",
            "Same Simulink model that ships to the real car (HIL parity)",
            "Co-simulation over IPC (deterministic stepping)",
            "Why it matters more than CARLA for our use case",
        ],
        code_anchor="src/scenic/simulators/dspace/simulator.py:82-147",
        speaker_note="",
    ),
    Slide(
        n=35, section=SECTION_T3,
        title="dSPACE / Simulink / Real Racecar -- The Toolchain",
        bullets=[
            "Simulink (MathWorks): graphical block-based vehicle dynamics model",
            "dSPACE: tooling stack on top of Simulink for automotive testing",
            "  ModelDesk -- configures the model (roads, traffic, scenario params)",
            "  ControlDesk -- runs experiments, displays variables",
            "  VEOS -- pure-software simulator running the SAME Simulink model",
            "Real racecar: ECU runs compiled code generated from the same model",
        ],
        code_anchor="src/scenic/simulators/dspace/simulator.py:82-147",
        speaker_note="",
    ),
    Slide(
        n=36, section=SECTION_T3,
        title="DSpaceSimulator + SyncStepBridge -- The Cosim Layer",
        bullets=[
            "DSpaceSimulator.__init__  -- simulator.py:82-147",
            "createSimulation()  -- simulator.py:237",
            "SyncStepBridge  -- python_listener/sync_step_bridge.py:9-168",
            "IPC client (VeosCoSimTestClientIpc.exe) spawned at simulator.py:196-202",
            "Flow: Scenic process -> SyncStepBridge -> IPC client -> VEOS",
        ],
        code_anchor="src/scenic/simulators/dspace/simulator.py:237",
        speaker_note="",
    ),
    Slide(
        n=37, section=SECTION_T3,
        title="MAPort: How State Crosses the Bridge",
        bullets=[
            "MAPort = dSPACE's variable-access API (read/write Simulink signals)",
            "One MAPort variable per car / sensor / actuator",
            "Mapping is configuration-driven (no codegen)",
            "Latency dominated by VEOS step time, NOT bridge overhead",
        ],
        code_anchor="src/scenic/simulators/dspace/simulator.py:544-686",
        speaker_note="",
    ),
    Slide(
        n=38, section=SECTION_T3,
        title="The `_racing_st_offset` Specifier",
        bullets=[
            "Parser: placement.py:26-58  (parse `(s_offset, t_offset)` from a tuple)",
            "Resolution: placement.py:409-428  (reads ego's s,t at sim init, applies offset)",
            "Gives `(s, t)` placement against ego's TTL frame",
            "Lives outside Scenic core because Scenic has no Frenet primitive",
            "Works, but is dSPACE-only -- not portable to other simulators",
        ],
        code_anchor="src/scenic/simulators/dspace/modeldesk/placement.py:409-428",
        missing_marker="Missing from Scenic: Frenet specifier should be domain-agnostic.",
        speaker_note="",
    ),

    # ---- Closing -- Feedback (divider + 5 slides) ----
    Slide(
        n=39, section=SECTION_FB, is_divider=True,
        title="Feedback",
        bullets=["Two concrete asks for Scenic core"],
        speaker_note="",
    ),
    Slide(
        n=40, section=SECTION_FB,
        title="Constructive Feedback Frame",
        bullets=[
            "Scenic's class / region / behavior model held up under racing pressure",
            "Two specific gaps surfaced: Frenet frames, state-machine behaviors",
            "Both are concrete proposals with API sketches and an implementation path",
            "Both come from real workarounds we wrote -- not speculation",
            "Neither is Scenic-was-wrong. Both are Scenic-could-grow.",
        ],
        speaker_note="",
    ),
    Slide(
        n=41, section=SECTION_FB,
        title="Ask 1: First-Class Frenet Frame",
        bullets=[
            "Workaround today: `_racing_st_offset` in placement.py:26-58 (~30 LOC, dSPACE-only)",
            "Confirmed not in Scenic core: zero matches for 'Frenet' / 's_offset' / 't_offset'",
            "Existing partial: LinearElement.flowFrom(point, distance) at roads.py:344 -- s only, no t",
            "Claim: Scenic core should expose `FrenetRegion(curve, s_range, t_range)` + `frenetOffset (s, t) from <obj>` specifier",
            "API sketch: `region = FrenetRegion(track.optimal_line, s=(0,100), t=(-3,3))`  and  `opp = new RacingCar with frenetOffset (gap_m, 5) from ego`",
            "Implementation: extend Polyline with frenet_at(s, t) (~30 LOC of projection math); specifier desugars to (x, y) at sample time",
            "Generalizes: any path-following domain (racing, ADAS lane-change, pedestrian crosswalks)",
        ],
        code_anchor="src/scenic/simulators/dspace/modeldesk/placement.py:26-58",
        speaker_note="",
    ),
    Slide(
        n=42, section=SECTION_FB,
        title="Ask 1: Same Logic, Side-by-Side",
        code_left=(
            "Today  --  dSPACE-only specifier",
            "# in S1_falsify.scenic:\n"
            "opp = new RacingCar with \\\n"
            "    _racing_st_offset (gap_m, 5)\n"
            "\n"
            "# placement.py:26-58 (dSPACE-specific)\n"
            "def _racing_st_offset_to_deltas(spec):\n"
            "    a, b = spec\n"
            "    if a == 'ahead': return (float(b), 0.0)\n"
            "    if a == 'left':  return (0.0, float(b))\n"
            "    # ... 'behind', 'right', numeric tuple\n"
            "\n"
            "# placement.py:409-428 resolution:\n"
            "#   reads ego._route_s_t at sim start,\n"
            "#   adds delta_s, delta_t,\n"
            "#   projects back to (x, y) via ego's\n"
            "#   route polyline.\n"
            "# Lives outside Scenic core; ~30 LOC;\n"
            "# dSPACE-only -- not portable."
        ),
        code_right=(
            "Proposed  --  Scenic core primitive",
            "# in any .scenic file:\n"
            "region = FrenetRegion(\n"
            "    track.optimal_line,\n"
            "    s=(0, 100), t=(-3, 3))\n"
            "opp = new RacingCar with \\\n"
            "    frenetOffset (gap_m, 5) from ego\n"
            "\n"
            "# scenic/core/regions.py\n"
            "class FrenetRegion(Region):\n"
            "    def __init__(curve, s_range, t_range):\n"
            "        ...\n"
            "    def frenet_at(s, t):\n"
            "        x, y, h = curve.point_at(s)\n"
            "        return Vector(x - t*sin(h),\n"
            "                      y + t*cos(h))\n"
            "\n"
            "# Specifier desugars at sample time -> (x, y).\n"
            "# Portable across all simulators."
        ),
        code_anchor="src/scenic/simulators/dspace/modeldesk/placement.py:26-58",
        speaker_note="",
    ),
    Slide(
        n=43, section=SECTION_FB,
        title="Ask 2: First-Class State-Machine Behaviors",
        bullets=[
            "Workaround today: state stored as `mode: str` on TacticalPlannerState dataclass (tactical_planner.py:330)",
            "Transitions: chained if/elif on `state.mode` inside `tactical_planner_step_v1()`; entry actions inlined ad-hoc beside every `state.mode = X` assignment",
            "Confirmed not in Scenic core: no `state` keyword in src/scenic/syntax/",
            "Existing primitive: `try-interrupt` (compiler.py:654-686) -- async transitions, but state must be hand-managed",
            "Claim: expose `state X:` blocks inside `behavior` definitions, with `on entry`, `on exit`, `when (cond): goto Y`",
            "Implementation: parser desugars `state` blocks into nested `try-interrupt` chains with a state variable; backwards-compatible (omit `state` = legacy behavior)",
            "Generalizes: every ADAS / autonomous-driving paper rolls its own FSM (CARLA examples, LGSVL navigators, our planner)",
        ],
        code_anchor="src/scenic/domains/racing/tactical_planner.py:330-396",
        speaker_note="",
    ),

    Slide(
        n=44, section=SECTION_FB,
        title="Ask 2: Same Logic, Side-by-Side",
        code_left=(
            "Today  --  tactical_planner.py",
            "if pick_left:\n"
            "    state.mode = COMMIT_PASS_LEFT\n"
            "    state.commit.side = 'left'\n"
            "    state.commit.start_s = sim_time_s\n"
            "    state.commit.until_s = sim_time_s + commit_hold_s\n"
            "    state.commit.candidate_count = ce_cycles\n"
            "    state.lateral_path_lock_side = 'left'\n"
            "    # ... 5 more entry-action lines\n"
            "elif pick_right:\n"
            "    state.mode = COMMIT_PASS_RIGHT\n"
            "    # ... mirror image\n"
            "elif state.mode == COMMIT_PASS_LEFT and lateral_clear:\n"
            "    state.mode = HOLD_PASS_LEFT\n"
            "    state.commit.hold_entry_s = sim_time_s\n"
            "    # ... mirror for COMMIT_PASS_RIGHT\n"
            "elif state.mode == HOLD_PASS_LEFT and merge_safe:\n"
            "    state.mode = FREE_RUN\n"
            "    # ... abort transitions, emergency overrides ..."
        ),
        code_right=(
            "Proposed  --  Scenic state-block syntax",
            "behavior TacticalPlanner():\n"
            "    state FREE_RUN:\n"
            "        when (strategy == 'pass_left'):\n"
            "            goto COMMIT_PASS_LEFT\n"
            "        when (strategy == 'pass_right'):\n"
            "            goto COMMIT_PASS_RIGHT\n"
            "\n"
            "    state COMMIT_PASS_LEFT:\n"
            "        on entry:\n"
            "            state.commit.side = 'left'\n"
            "            state.commit.start_s = sim_time\n"
            "            take SetTtl('left')\n"
            "        when (lateral_clear):\n"
            "            goto HOLD_PASS_LEFT\n"
            "        when (path_collision_predicted):\n"
            "            goto ABORT_PASS\n"
            "\n"
            "    state HOLD_PASS_LEFT:\n"
            "        when (merge_safe): goto FREE_RUN"
        ),
        code_anchor="src/scenic/domains/racing/tactical_planner.py:330-396",
        speaker_note="",
    ),

    # ---- Close (36-37) ----
    Slide(
        n=45, section=SECTION_CLOSE,
        title="Recap: 3 Topics, 2 Asks",
        bullets=[
            "Topic 1: scenario banks (F-bank + falsifiable)",
            "Topic 2: racing library (RacingCar + MPC behavior + tactical planner + falsifier loop)",
            "Topic 3: dSPACE bridge (light architectural sketch)",
            "Asks: First-class Frenet frame; first-class state-machine behaviors",
        ],
        speaker_note="",
    ),
    Slide(
        n=46, section=SECTION_CLOSE, is_divider=True,
        title="Thanks For Listening",
        bullets=["Open the floor for questions"],
        speaker_note="",
    ),
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _set_run(run, text, font, size, color, bold=False, italic=False):
    run.text = text
    run.font.name = font
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def _add_textbox(slide, box, font, size, color, text, bold=False, italic=False, align=None):
    left, top, width, height = box
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    _set_run(run, text, font, size, color, bold=bold, italic=italic)
    return tb, run


def _add_code_anchor_with_link(slide, text, url):
    """Render the code-anchor line, hyperlinking the whole text to `url`.

    Whole-line hyperlink keeps the visual simple: the entire row is
    clickable. Falls back to plain gray text when no URL is available.
    """
    left, top, width, height = CODE_BOX
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    label = "> " + text
    if url:
        _set_run(run, label, CODE_FONT, CODE_SIZE, COLOR_LINK)
        run.hyperlink.address = url
    else:
        _set_run(run, label, CODE_FONT, CODE_SIZE, COLOR_CODE)


def _add_labeled_code_block(slide, box, label, code):
    """Render one label header + a monospace code box below it."""
    left, top, width, height = box
    label_h = Inches(0.40)
    # Label
    lb = slide.shapes.add_textbox(left, top, width, label_h)
    lp = lb.text_frame.paragraphs[0]
    lp.alignment = PP_ALIGN.LEFT
    lr = lp.add_run()
    _set_run(lr, label, BODY_FONT, Pt(14), COLOR_TITLE, bold=True)
    # Code panel
    cb = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + label_h, width, height - label_h
    )
    cb.fill.solid()
    cb.fill.fore_color.rgb = COLOR_VIDEO_FILL
    cb.line.color.rgb = COLOR_VIDEO_BORDER
    cb.line.width = Pt(0.75)
    tf = cb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    lines = code.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        _set_run(run, line if line else " ", CODE_FONT, Pt(10.5), COLOR_BODY)


def _add_code_block_pair(slide, left_pair, right_pair):
    """Render two side-by-side labeled code blocks for the comparison layout."""
    lbox = (Inches(0.5), Inches(1.3), Inches(6.15), Inches(4.7))
    rbox = (Inches(6.7), Inches(1.3), Inches(6.15), Inches(4.7))
    _add_labeled_code_block(slide, lbox, left_pair[0], left_pair[1])
    _add_labeled_code_block(slide, rbox, right_pair[0], right_pair[1])


def _add_bullets(slide, bullets):
    left, top, width, height = BODY_BOX
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        _set_run(run, "•  " + bullet, BODY_FONT, BODY_SIZE, COLOR_BODY)


def _add_divider_layout(slide, title_text, subtitle_lines):
    """Topic-divider layout: huge centered title + small subtitle line(s).

    `subtitle_lines` is a list of strings; each renders as a separate
    centered paragraph below the title at progressively smaller weight.
    First line at Pt(22) italic; subsequent lines at Pt(16) regular.

    Replaces the normal title-at-top + bullets-below layout. Used to
    visually break sections so the audience sees a clean transition
    between major topics during the talk, and for the cover/thanks
    book-end slides.
    """
    box_top = Inches(2.6)
    title_h = Inches(1.4)
    sub_h = Inches(1.4)

    # Big centered title
    tb = slide.shapes.add_textbox(Inches(0.5), box_top, Inches(12.3), title_h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, title_text, TITLE_FONT, Pt(56), COLOR_TITLE, bold=True)

    # Subtitle line(s): first at 22pt italic, rest at 16pt
    if subtitle_lines:
        sb = slide.shapes.add_textbox(
            Inches(0.5), box_top + title_h + Inches(0.2),
            Inches(12.3), sub_h
        )
        sf = sb.text_frame
        sf.word_wrap = True
        for i, line in enumerate(subtitle_lines):
            para = sf.paragraphs[0] if i == 0 else sf.add_paragraph()
            para.alignment = PP_ALIGN.CENTER
            srun = para.add_run()
            if i == 0:
                _set_run(srun, line, BODY_FONT, Pt(22), COLOR_FOOTER, italic=True)
            else:
                _set_run(srun, line, BODY_FONT, Pt(16), COLOR_FOOTER, italic=False)


def _add_state_box(slide, x, y, w, h, label):
    """Rounded rectangle node for an FSM state, with centered bold text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE5, 0xED, 0xFA)
    shape.line.color.rgb = COLOR_TITLE
    shape.line.width = Pt(1.0)
    tf = shape.text_frame
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, label, BODY_FONT, Pt(10), COLOR_BODY, bold=True)
    return shape


def _add_arrow(slide, x1, y1, x2, y2, color=None, width=None):
    """Straight line connector with an arrowhead at the end (x2, y2)."""
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn
    from lxml import etree

    # Coerce to int EMU; arithmetic on Inches values can produce floats
    # (e.g. Inches(1.55) / 2 = 708660.0), and the OOXML schema requires
    # integer EMU values on <a:off x="..." y="..."/>. Without this the
    # generated .pptx fails to round-trip through python-pptx readers.
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2)
    )
    line = connector.line
    line.color.rgb = color if color is not None else COLOR_BODY
    line.width = width if width is not None else Pt(1.0)
    # Arrowhead at tail (i.e. the destination end) via direct XML
    line_xml = line._get_or_add_ln()
    tailEnd = etree.SubElement(line_xml, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("h", "med")
    return connector


def _add_label(slide, cx, cy, text, font_size=Pt(8.5)):
    """Tiny italic label centered at (cx, cy). Used for FSM transition tags."""
    w = Inches(1.2)
    h = Inches(0.22)
    tb = slide.shapes.add_textbox(cx - w / 2, cy - h / 2, w, h)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, text, BODY_FONT, font_size, COLOR_FOOTER, italic=True)
    return tb


def _add_fsm_diagram(slide):
    """Render the planner state machine in the right panel of the slide.

    5 logical states (COMMIT/HOLD have left/right variants -- collapsed
    here for diagram clarity). Returns to FREE_RUN are rendered as long
    diagonal arrows around the periphery to keep the topology readable.
    """
    # Diagram panel: x in [9.0, 13.2], y in [1.4, 6.1]
    w_node = Inches(1.55)
    h_node = Inches(0.45)

    # Header label
    _add_label(slide, Inches(11.1), Inches(1.32),
               "Planner state machine",
               font_size=Pt(10))

    # Node centers (chosen by hand)
    free_cx, free_cy = Inches(11.1), Inches(1.85)
    foll_cx, foll_cy = Inches(11.1), Inches(2.75)
    commit_cx, commit_cy = Inches(10.05), Inches(3.85)
    abort_cx, abort_cy = Inches(12.15), Inches(3.85)
    hold_cx, hold_cy = Inches(10.05), Inches(4.95)

    # Place the 5 state boxes
    half_w = w_node / 2
    half_h = h_node / 2
    for label, cx, cy in [
        ("FREE_RUN", free_cx, free_cy),
        ("FOLLOW", foll_cx, foll_cy),
        ("COMMIT_PASS  L/R", commit_cx, commit_cy),
        ("ABORT_PASS", abort_cx, abort_cy),
        ("HOLD_PASS  L/R", hold_cx, hold_cy),
    ]:
        _add_state_box(slide, cx - half_w, cy - half_h, w_node, h_node, label)

    # Forward arrows (left side of vertical pairs to leave room for return-arrows)
    off = Inches(0.10)  # small horizontal offset to separate up/down arrows

    # FREE_RUN -> FOLLOW
    _add_arrow(slide,
               free_cx - off, free_cy + half_h,
               foll_cx - off, foll_cy - half_h)
    # FOLLOW -> FREE_RUN  (return)
    _add_arrow(slide,
               foll_cx + off, foll_cy - half_h,
               free_cx + off, free_cy + half_h)
    # FOLLOW -> COMMIT_PASS  (diagonal down-left)
    _add_arrow(slide,
               foll_cx - half_w / 3, foll_cy + half_h,
               commit_cx + half_w / 3, commit_cy - half_h)
    # COMMIT_PASS -> ABORT_PASS  (rightward)
    _add_arrow(slide,
               commit_cx + half_w, commit_cy,
               abort_cx - half_w, abort_cy)
    # COMMIT_PASS -> HOLD_PASS  (down)
    _add_arrow(slide,
               commit_cx, commit_cy + half_h,
               hold_cx, hold_cy - half_h)
    # HOLD_PASS -> FREE_RUN  (long arrow up the LEFT side, around the diagram)
    _add_arrow(slide,
               hold_cx - half_w, hold_cy,
               Inches(9.15), Inches(1.85))
    _add_arrow(slide,
               Inches(9.15), Inches(1.85),
               free_cx - half_w, free_cy)
    # ABORT_PASS -> FREE_RUN  (long arrow up the RIGHT side)
    _add_arrow(slide,
               abort_cx + half_w, abort_cy,
               Inches(13.05), Inches(1.85))
    _add_arrow(slide,
               Inches(13.05), Inches(1.85),
               free_cx + half_w, free_cy)

    # Transition annotations (italic gray, near each arrow's midpoint).
    # Placed on the OUTSIDE of the diagram where possible to avoid
    # collisions with other arrows.
    _add_label(slide, Inches(10.55), Inches(2.20), "fellow ahead", font_size=Pt(8))
    _add_label(slide, Inches(11.65), Inches(2.20), "clear",         font_size=Pt(8))
    _add_label(slide, Inches(10.65), Inches(3.15), "strategy=pass", font_size=Pt(8))
    _add_label(slide, Inches(10.85), Inches(3.72), "collision risk",font_size=Pt(8))
    _add_label(slide, Inches( 9.45), Inches(4.30), "lateral clear", font_size=Pt(8))
    _add_label(slide, Inches( 8.95), Inches(3.40), "merge safe",    font_size=Pt(8))
    _add_label(slide, Inches(13.25), Inches(2.85), "recovered",     font_size=Pt(8))


def _add_video_placeholder(slide, spec):
    left, top, width, height = VIDEO_BOX
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_VIDEO_FILL
    shape.line.color.rgb = COLOR_VIDEO_BORDER
    shape.line.width = Pt(1.0)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, "VIDEO PLACEHOLDER", BODY_FONT, Pt(14), COLOR_FOOTER, bold=True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    run2 = p2.add_run()
    _set_run(run2, spec, BODY_FONT, Pt(13), COLOR_BODY, italic=True)


def _add_footer(slide, slide_n, total, section):
    left, top, width, height = FOOTER_BOX
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    _set_run(run, f"{section}", BODY_FONT, FOOTER_SIZE, COLOR_FOOTER, italic=True)
    rb = slide.shapes.add_textbox(Inches(11.5), top, Inches(1.7), height)
    rtf = rb.text_frame
    rp = rtf.paragraphs[0]
    rp.alignment = PP_ALIGN.RIGHT
    rrun = rp.add_run()
    _set_run(rrun, f"{slide_n} / {total}", BODY_FONT, FOOTER_SIZE, COLOR_FOOTER)


def _set_speaker_note(slide, note):
    if not note:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = ""
    p = tf.paragraphs[0]
    run = p.add_run()
    _set_run(run, note, BODY_FONT, NOTE_SIZE, COLOR_BODY)


# ---------------------------------------------------------------------------
# Render
# ---------------------------------------------------------------------------

def build_deck(out_path: Path) -> int:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]

    total = len(SLIDES)
    for sd in SLIDES:
        slide = prs.slides.add_slide(blank_layout)

        if sd.is_divider:
            # Topic-divider layout: skip the small top title and the
            # body box; render a huge centered title + subtitle line(s)
            # instead. `bullets` becomes the list of subtitle lines.
            _add_divider_layout(slide, sd.title, list(sd.bullets))
            _add_footer(slide, sd.n, total, sd.section)
            _set_speaker_note(slide, sd.speaker_note)
            continue

        _add_textbox(slide, TITLE_BOX, TITLE_FONT, TITLE_SIZE, COLOR_TITLE,
                     sd.title, bold=True)

        if sd.code_left and sd.code_right:
            # Side-by-side code-comparison layout. Replaces the bullets body.
            _add_code_block_pair(slide, sd.code_left, sd.code_right)
        elif sd.bullets:
            _add_bullets(slide, sd.bullets)

        if sd.video_spec:
            _add_video_placeholder(slide, sd.video_spec)

        if sd.fsm_diagram:
            _add_fsm_diagram(slide)

        if sd.code_anchor:
            url = _anchor_url(sd.code_anchor)
            _add_code_anchor_with_link(slide, sd.code_anchor, url)

        if sd.missing_marker:
            _add_textbox(slide, MISSING_BOX, BODY_FONT, MISSING_SIZE, COLOR_MISSING,
                         sd.missing_marker, italic=True)

        _add_footer(slide, sd.n, total, sd.section)

        _set_speaker_note(slide, sd.speaker_note)

    prs.save(str(out_path))
    return total


def main() -> int:
    out_path = Path(__file__).parent / "scenic_racing_review.pptx"
    n = build_deck(out_path)
    print(f"[deck] generated: {out_path}")
    print(f"[deck] slides:    {n}")
    print(f"[deck] sections:  open=4, t1=5, t2a=7, t2b=9, t2c=7, t3=6, fb=6, close=2")
    print(f"[deck] dividers:  slide 1 (cover), 5/10/17/26/33/39 (topics), 46 (thanks)")
    print(f"[deck] markers:   slides 15, 23, 38 carry 'Missing-from-Scenic' lines")
    print(f"[deck] videos:    slides 2, 7, 8, 30 reserved for video placeholders")
    print(f"[deck] fsm:       slide 23 has the planner state-machine diagram")
    print(f"[deck] code-cmp:  slides 42 (Ask 1) and 44 (Ask 2) have side-by-side code")
    print(f"[deck] notes:     speaker notes intentionally blank (talk delivered)")
    print(f"[deck] hyperlinks: code-anchor lines link to {GITHUB_BASE}/tree/{BRANCH}/...")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
